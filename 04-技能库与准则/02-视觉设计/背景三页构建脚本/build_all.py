#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
背景三页-新版.pptx 完整生成脚本（第一/二/三页）
依据：技能库&准则/03-专家团与技能路由/02-背景页改造工作法.md
"""
import sys, os
sys.path.insert(0, os.path.dirname(__file__))
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import design_system as ds
from design_system import IN, WHITE, GOLD, CYAN, GRAY, PANEL_FILL, PANEL_LINE, RED_WARN, BADGE_FILL, BOTTOM_BAR_FILL

ASSETS = "/home/user/sucheng/assets"
BUILD = "/home/user/sucheng/build_assets"

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)


def arrow_down(slide, x, y, w=0.22, h=0.16, color=CYAN):
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Emu(int(x)), Emu(int(y)), Emu(int(w * IN)), Emu(int(h * IN)))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = color; arrow.line.fill.background(); arrow.shadow.inherit = False
    return arrow


def arrow_right(slide, x, y, w=0.28, h=0.2, color=GOLD):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(int(x)), Emu(int(y)), Emu(int(w * IN)), Emu(int(h * IN)))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = color; arrow.line.fill.background(); arrow.shadow.inherit = False
    return arrow


# =====================================================================
# 第一页：PEEK 基复合材料应用与 LPBF 成形工艺
# =====================================================================
s1 = ds.new_slide(prs)
ds.add_header(s1, prs, "PEEK 基复合材料应用需求明确，倒逼高温 LPBF 成形工艺升级")

LEFT_X = 0.4 * IN
LEFT_W = 6.55 * IN
RIGHT_X = 7.15 * IN
RIGHT_W = 5.2 * IN
TOP_Y = 1.12 * IN
TOP_H = 3.05 * IN

ds.add_panel(s1, LEFT_X, TOP_Y, LEFT_W, TOP_H, header="应用场景 —— PEEK 基复合材料已渗透关键运动部位")
apps = [
    ("icon_hinge.png", "航天器铰链", "卫星太阳翼展开机构、锁钩机构"),
    ("icon_skull.png", "PEEK/GF 头骨植入物", "玻纤增强，生物相容+承力兼顾"),
    ("icon_bracket.png", "PEEK/CF 承力支架", "碳纤增强，轻量化替代金属"),
    ("icon_bearing.png", "PEEK/PTFE 低摩擦构件", "自润滑轴套，免维护运行"),
]
grid_top = TOP_Y + 0.42 * IN
cell_w = (LEFT_W - 0.3 * IN) / 2
cell_h = (TOP_H - 0.55 * IN) / 2
for idx, (icon, name, desc) in enumerate(apps):
    col = idx % 2
    row = idx // 2
    cx = LEFT_X + 0.15 * IN + col * cell_w
    cy = grid_top + row * (cell_h + 0.06 * IN)
    icon_size = cell_h * 0.62
    ds.add_picture_framed(s1, f"{ASSETS}/{icon}", cx + 0.08 * IN, cy + 0.05 * IN, icon_size, icon_size)
    ds.add_text(s1, cx + icon_size + 0.2 * IN, cy + 0.08 * IN, cell_w - icon_size - 0.35 * IN, 0.35 * IN,
                name, size=13.5, color=WHITE, bold=True)
    ds.add_text(s1, cx + icon_size + 0.2 * IN, cy + 0.42 * IN, cell_w - icon_size - 0.35 * IN, cell_h - 0.5 * IN,
                desc, size=10.5, color=GRAY)

ds.add_panel(s1, RIGHT_X, TOP_Y, RIGHT_W, TOP_H, header="传统加工 → 普通增材 → 高温 LPBF 的能力递进")
steps = [
    ("icon_cnc.png", "传统加工受限", "铣削/线切割难以成形复杂内腔与一体化流道结构"),
    ("icon_fdm.png", "普通增材难适配", "PEEK 熔点 343℃，FDM/普通 SLS 精度低、易翘曲开裂"),
    ("icon_lpbf.png", "高温 LPBF 可成形复杂结构", "高温预热 + 精密光路，一体化打印点阵/流道等复杂构型"),
]
step_top = TOP_Y + 0.38 * IN
step_h = (TOP_H - 0.5 * IN) / 3
for i, (icon, name, desc) in enumerate(steps):
    cy = step_top + i * (step_h + 0.02 * IN)
    icon_size = step_h * 0.72
    ds.add_picture_framed(s1, f"{ASSETS}/{icon}", RIGHT_X + 0.15 * IN, cy + (step_h - icon_size) / 2, icon_size, icon_size)
    tx = RIGHT_X + 0.15 * IN + icon_size + 0.18 * IN
    tw = RIGHT_W - (icon_size + 0.5 * IN)
    name_color = GOLD if i == 2 else WHITE
    ds.add_text(s1, tx, cy + 0.03 * IN, tw, 0.3 * IN, f"{i+1}. {name}", size=13, color=name_color, bold=True)
    ds.add_text(s1, tx, cy + 0.34 * IN, tw, step_h - 0.4 * IN, desc, size=10, color=GRAY)
    if i < 2:
        arrow_down(s1, RIGHT_X + RIGHT_W - 0.35 * IN, cy + step_h - 0.02 * IN)

BOT_Y = TOP_Y + TOP_H + 0.12 * IN
BOT_H = 1.85 * IN
ds.add_panel(s1, LEFT_X, BOT_Y, LEFT_W, BOT_H, header="LPBF 成形原理与 PEEK 复合材料制件")
proc = [
    ("icon", "icon_lpbf.png", "① 成形过程", "激光逐层扫描粉床，逐层熔融堆积"),
    ("build", "peek_powder_white.png", "② PEEK 基复合粉末", "自润滑填料+PEEK基体，按需定制配比"),
    ("build", "peek_powder_black.png", "③ 成形构件", "一体化复杂结构，减重且免维护"),
]
pw = (LEFT_W - 0.3 * IN) / 3
for i, (kind, img, name, desc) in enumerate(proc):
    px = LEFT_X + 0.15 * IN + i * pw
    py = BOT_Y + 0.42 * IN
    isize = 0.85 * IN
    path = f"{ASSETS}/{img}" if kind == "icon" else f"{BUILD}/{img}"
    ds.add_picture_framed(s1, path, px + (pw - isize) / 2, py, isize, isize)
    ds.add_text(s1, px, py + isize + 0.06 * IN, pw, 0.28 * IN, name, size=11.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    ds.add_text(s1, px + 0.05 * IN, py + isize + 0.32 * IN, pw - 0.1 * IN, 0.5 * IN, desc, size=8.5, color=GRAY, align=PP_ALIGN.CENTER)

ds.add_panel(s1, RIGHT_X, BOT_Y, RIGHT_W, BOT_H, header="传统材料/工艺性能对比（国内外头部企业口径）")
table_rows = [
    ("拉伸强度 (MPa)", "95~105", "86~90"),
    ("材料摩擦系数", "0.18~0.22", "0.28~0.35"),
    ("成形尺寸误差", "±1.5%~±3%", "±2%~±4%"),
]
th = (BOT_H - 0.75 * IN) / 3
tw0 = RIGHT_W - 0.3 * IN
col1, col2, col3 = 0.42, 0.29, 0.29
ty0 = BOT_Y + 0.42 * IN
hx = RIGHT_X + 0.15 * IN
ds.add_text(s1, hx, ty0, tw0 * col1, 0.28 * IN, "指标", size=10.5, color=CYAN, bold=True)
ds.add_text(s1, hx + tw0 * col1, ty0, tw0 * col2, 0.28 * IN, "国外", size=10.5, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
ds.add_text(s1, hx + tw0 * (col1 + col2), ty0, tw0 * col3, 0.28 * IN, "国产", size=10.5, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
for i, (metric, out_v, cn_v) in enumerate(table_rows):
    ry = ty0 + 0.3 * IN + i * th
    ds.add_text(s1, hx, ry, tw0 * col1, th, metric, size=10, color=WHITE)
    ds.add_text(s1, hx + tw0 * col1, ry, tw0 * col2, th, out_v, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    ds.add_text(s1, hx + tw0 * (col1 + col2), ry, tw0 * col3, th, cn_v, size=10, color=GRAY, align=PP_ALIGN.CENTER)
ds.add_text(s1, hx, ty0 + 0.3 * IN + 3 * th + 0.02 * IN, tw0, 0.3 * IN,
            "注：比较国内外头部企业产品性能", size=8, color=GRAY)

ds.add_bottom_bar(s1, prs, [
    ("PEEK 基复合材料应用需求明确，但", WHITE, False),
    ("高熔点、低导热和温度窗口敏感", GOLD, True),
    ("，对 LPBF 成形设备提出更高要求。", WHITE, False),
])

print("slide 1 built")

# =====================================================================
# 第二页：国产高温 LPBF 成形装备的能力不足
# =====================================================================
s2 = ds.new_slide(prs)
ds.add_header(s2, prs, "国产高温 LPBF 成形装备能力不足，过程温控与稳定成形受制于人")

L2X = 0.4 * IN
L2W = 6.7 * IN
R2X = 7.3 * IN
R2W = 5.05 * IN
T2Y = 1.12 * IN

# 左上：装备能力五维对比
CAP_H = 2.15 * IN
EVI_H_FIXED = 3.19 * IN
ds.add_panel(s2, L2X, T2Y, L2W, CAP_H, header="装备能力对比 —— 不止比材料，更要比装备")
caps = [
    ("最高预热温度", "385℃（德国 EOS P800）", "340~380℃（国产装备，口径待统一）"),
    ("温度场均匀性", "8 区独立温控，波动小", "分区控温程度低，热累积明显"),
    ("过程监测能力", "在线激光功率反馈调节", "缺乏实时多模态测温反馈"),
    ("成形稳定性/精度", "翘曲概率低，尺寸误差小", "翘曲、结块、粉层开裂时有发生"),
]
cap_h = (CAP_H - 0.45 * IN) / 4
cy0 = T2Y + 0.42 * IN
cw0 = L2W - 0.3 * IN
c1, c2, c3 = 0.24, 0.38, 0.38
for i, (metric, out_v, cn_v) in enumerate(caps):
    ry = cy0 + i * cap_h
    cx = L2X + 0.15 * IN
    ds.add_text(s2, cx, ry, cw0 * c1, cap_h, metric, size=10, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    ds.add_text(s2, cx + cw0 * c1, ry, cw0 * c2, cap_h, out_v, size=9, color=CYAN, anchor=MSO_ANCHOR.MIDDLE)
    ds.add_text(s2, cx + cw0 * (c1 + c2), ry, cw0 * c3, cap_h, cn_v, size=9, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)

# 左下：现象证据图
EVI_Y = T2Y + CAP_H + 0.12 * IN
EVI_H = EVI_H_FIXED
ds.add_panel(s2, L2X, EVI_Y, L2W, EVI_H, header="现象证据 —— 当前工艺仍存在热累积、结块与翘曲")
ev_w = (L2W - 0.35 * IN) / 2
ev_img_h = EVI_H - 0.85 * IN
ds.add_picture_framed(s2, f"{BUILD}/warp_evidence_2.jpg", L2X + 0.15 * IN, EVI_Y + 0.42 * IN, ev_w, ev_img_h)
ds.add_text(s2, L2X + 0.15 * IN, EVI_Y + 0.42 * IN + ev_img_h + 0.05 * IN, ev_w, 0.3 * IN,
            "热累积现象加剧却无法实时控温", size=10, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
ds.add_picture_framed(s2, f"{BUILD}/warp_evidence_1.jpg", L2X + 0.25 * IN + ev_w, EVI_Y + 0.42 * IN, ev_w, ev_img_h)
ds.add_text(s2, L2X + 0.25 * IN + ev_w, EVI_Y + 0.42 * IN + ev_img_h + 0.05 * IN, ev_w, 0.3 * IN,
            "翘曲现象 / 粉层开裂", size=10, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# 右上：政策
POL_H = 2.7 * IN
ds.add_panel(s2, R2X, T2Y, R2W, POL_H, header="政策导向 —— 国产替代需求明确")
img_w = 2.05 * IN
ds.add_picture_framed(s2, f"{BUILD}/policy_meeting.png", R2X + 0.15 * IN, T2Y + 0.42 * IN, img_w, img_w * 0.56)
ds.add_picture_framed(s2, f"{BUILD}/policy_text.png", R2X + 0.15 * IN, T2Y + 0.42 * IN + img_w * 0.56 + 0.08 * IN,
                       R2W - 0.3 * IN, (R2W - 0.3 * IN) * (194 / 1509))
ds.add_text(s2, R2X + 0.15 * IN, T2Y + POL_H - 0.4 * IN, R2W - 0.3 * IN, 0.38 * IN,
            "“十五五”规划建议：加快新能源、新材料、航空航天等战略性新兴产业集群发展",
            size=9, color=GRAY)

# 右下：市场规模 + 依赖表述
MKT_Y = T2Y + POL_H + 0.12 * IN
MKT_H = 2.64 * IN
ds.add_panel(s2, R2X, MKT_Y, R2W, MKT_H, header="全球市场规模潜力巨大")
mw = (R2W - 0.35 * IN) / 2
mimg_h = 1.35 * IN
ds.add_picture_framed(s2, f"{BUILD}/chart_lpbf_market.png", R2X + 0.15 * IN, MKT_Y + 0.42 * IN, mw, mimg_h)
ds.add_text(s2, R2X + 0.15 * IN, MKT_Y + 0.42 * IN + mimg_h + 0.03 * IN, mw, 0.42 * IN,
            "全球 LPBF 技术市场规模\n（单位：亿元）", size=8.5, color=GRAY, align=PP_ALIGN.CENTER)
ds.add_picture_framed(s2, f"{BUILD}/chart_peek_market.png", R2X + 0.25 * IN + mw, MKT_Y + 0.42 * IN, mw, mimg_h)
ds.add_text(s2, R2X + 0.25 * IN + mw, MKT_Y + 0.42 * IN + mimg_h + 0.03 * IN, mw, 0.42 * IN,
            "全球 PEEK 粉末市场规模\n（单位：亿元）", size=8.5, color=GRAY, align=PP_ALIGN.CENTER)
ds.add_text(s2, R2X + 0.15 * IN, MKT_Y + MKT_H - 0.38 * IN, R2W - 0.3 * IN, 0.35 * IN,
            "高端复合粉末与专用装备高度依赖英国 Victrex、德国 Evonik 等企业",
            size=9, color=GOLD, bold=True)

ds.add_bottom_bar(s2, prs, [
    ("国产 PEEK 基复合粉末与专用高温 LPBF 装备协同不足，", WHITE, False),
    ("过程温控和稳定成形能力仍受制于人", GOLD, True),
    ("。", WHITE, False),
])

print("slide 2 built")

# =====================================================================
# 第三页：核心问题与解决方案
# =====================================================================
s3 = ds.new_slide(prs)
ds.add_header(s3, prs, "聚焦高温 LPBF 装备核心问题与解决路径")

T3Y = 1.15 * IN
P_H = 2.55 * IN
P_W = (12192000 - Emu(int(0.4 * IN)) * 2 - Emu(int(0.3 * IN)) * 2) / 3

problems = [
    ("icon_heat.png", "热场控制不足", "激光扫描瞬时高温与铺粉降温交替，热累积效应加剧，成形精度难以保证"),
    ("icon_powderbed.png", "粉床稳定性差", "铺粉不均、层间结合弱，PEEK 基复合粉末易结块、翘曲甚至粉层开裂"),
    ("icon_quality.png", "成形质量波动", "批次一致性差，复杂构件合格率低，难以满足工程化交付要求"),
]
for i, (icon, name, desc) in enumerate(problems):
    px = 0.4 * IN + i * (P_W + 0.3 * IN)
    ds.add_panel(s3, px, T3Y, P_W, P_H)
    isize = 1.0 * IN
    ds.add_picture_framed(s3, f"{ASSETS}/{icon}", px + (P_W - isize) / 2, T3Y + 0.18 * IN, isize, isize)
    ds.add_text(s3, px + 0.15 * IN, T3Y + isize + 0.28 * IN, P_W - 0.3 * IN, 0.35 * IN,
                name, size=15, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    ds.add_text(s3, px + 0.2 * IN, T3Y + isize + 0.68 * IN, P_W - 0.4 * IN, P_H - isize - 0.75 * IN,
                desc, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# 下半部分：解决路径
SOL_Y = T3Y + P_H + 0.35 * IN
SOL_H = 2.1 * IN
ds.add_panel(s3, 0.4 * IN, SOL_Y - 0.32 * IN, 12192000 - Emu(int(0.8 * IN)), SOL_H + 0.32 * IN,
             header="解决路径 —— 把缺陷控制前移到打印过程中")
sols = [
    ("①", "实时监测", "多模态点面协同测温，精确感知高温粉床热环境"),
    ("②", "精确调控", "多区域协同温度场闭环智能调控，抑制热累积与翘曲"),
    ("③", "数据采集", "构建工艺-微结构-性能数据库，沉淀多工况实验数据"),
    ("④", "映射关系", "机器学习建模指导工艺优化，反演最优成形参数组合"),
]
sw = (12192000 - Emu(int(0.8 * IN)) - Emu(int(0.3 * IN))) / 4
for i, (num, name, desc) in enumerate(sols):
    sx = 0.4 * IN + i * (sw + 0.1 * IN)
    sy = SOL_Y + 0.12 * IN
    circle = s3.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(sx + sw / 2 - 0.3 * IN)), Emu(int(sy)), Emu(int(0.6 * IN)), Emu(int(0.6 * IN)))
    circle.fill.solid(); circle.fill.fore_color.rgb = BADGE_FILL
    circle.line.color.rgb = CYAN; circle.line.width = Pt(1.25); circle.shadow.inherit = False
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = GOLD; r.font.name = ds.FONT
    ds.add_text(s3, sx, sy + 0.72 * IN, sw, 0.3 * IN, name, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    ds.add_text(s3, sx + 0.1 * IN, sy + 1.05 * IN, sw - 0.2 * IN, 0.85 * IN, desc, size=9, color=GRAY, align=PP_ALIGN.CENTER)
    if i < 3:
        arrow_right(s3, sx + sw - 0.02 * IN, sy + 0.18 * IN)

ds.add_bottom_bar(s3, prs, [
    ("项目聚焦高温 LPBF 成形设备，通过", WHITE, False),
    ("温度场监测和过程调控", GOLD, True),
    ("，把缺陷控制前移到打印过程中，提高 PEEK 基复合构件成形稳定性。", WHITE, False),
])

print("slide 3 built")

prs.save("/home/user/sucheng/背景三页-新版.pptx")
print("saved: 背景三页-新版.pptx")
