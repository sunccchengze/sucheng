"""
逐块dump PPTX源文件,显式用UTF-8写文件,文本字段尝试按gbk解码。
"""
from pptx import Presentation
import sys, io

# stdout 重写成 UTF-8
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
sys.stderr = sys.stdout

def emu_to_in(v):
    return v/914400.0

def safe(s):
    """尝试把乱码字符串按gbk恢复"""
    if s is None:
        return None
    if not isinstance(s, str):
        s = str(s)
    try:
        return s.encode("latin-1").decode("gbk", errors="replace")
    except Exception:
        return s

def iter_text_frame(tf, out_lines, indent=0):
    pad = "  " * indent
    for p in tf.paragraphs:
        for r in p.runs:
            txt = safe(r.text).replace("\n", "\\n")
            if not txt.strip():
                continue
            sz = r.font.size
            sz_pt = sz.pt if sz else None
            bold = r.font.bold
            color = None
            try:
                if r.font.color and r.font.color.rgb:
                    color = str(r.font.color.rgb)
            except Exception:
                pass
            font_name = r.font.name
            out_lines.append(
                f"{pad}run: sz={sz_pt} font={font_name} bold={bold} color={color} text={txt!r}"
            )

prs = Presentation("D:/sucheng/我现在的版本.pptx")
print(f"Total slides: {len(prs.slides)}", file=sys.stderr)
for i, slide in enumerate(prs.slides, 1):
    out = []
    out.append(f"\n{'='*70}\nSLIDE {i}\n{'='*70}")
    for shp in slide.shapes:
        name = safe(shp.name)
        left = emu_to_in(shp.left or 0)
        top  = emu_to_in(shp.top or 0)
        w    = emu_to_in(shp.width or 0)
        h    = emu_to_in(shp.height or 0)
        type_name = shp.shape_type
        out.append(f"\n-- shape: name={name!r} type={type_name}")
        out.append(f"   pos: left={left:.2f}in top={top:.2f}in  size: w={w:.2f}in h={h:.2f}in")
        if shp.has_text_frame:
            iter_text_frame(shp.text_frame, out, indent=1)
        if hasattr(shp, "has_table") and shp.has_table:
            try:
                tbl = shp.table
                for r_i, row in enumerate(tbl.rows):
                    cells = [safe(c.text).replace("\n"," ") for c in row.cells]
                    out.append(f"   table row {r_i}: {cells}")
            except Exception as e:
                out.append(f"   table read err: {e}")
    sys.stdout.write("\n".join(out))
    sys.stdout.write("\n")
sys.stdout.flush()
