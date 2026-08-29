#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
塑成非凡 · 背景三页新版 生成脚本
依据《背景修改建议-0811.docx》与 03-专家团与技能路由/02-背景页改造工作法.md 执行。
视觉美学官 + 信息设计官 联合把关：深蓝科技风、金/青高亮、发丝线分隔、Swiss Grid 对齐。
"""
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---------- 设计系统 ----------
BG_TOP = RGBColor(0x07, 0x0F, 0x22)
BG_BOTTOM = RGBColor(0x14, 0x30, 0x5C)
PANEL_FILL = RGBColor(0x11, 0x24, 0x46)
PANEL_LINE = RGBColor(0x3D, 0x7A, 0xB8)
CYAN = RGBColor(0x4FC3F7 >> 16 & 0xFF, 0x4FC3F7 >> 8 & 0xFF, 0x4FC3F7 & 0xFF)
GOLD = RGBColor(0xF2, 0xC4, 0x4C)
WHITE = RGBColor(0xF2, 0xF6, 0xFC)
GRAY = RGBColor(0xB8, 0xC6, 0xDA)
BADGE_FILL = RGBColor(0x1B, 0x3E, 0x74)
BOTTOM_BAR_FILL = RGBColor(0x0A, 0x18, 0x33)
RED_WARN = RGBColor(0xE0, 0x5A, 0x47)

FONT = "微软雅黑"

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


def set_gradient_bg(slide, prs):
    """整页深蓝科技渐变背景"""
    left = top = 0
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, prs.slide_width, prs.slide_height)
    shape.shadow.inherit = False
    shape.line.fill.background()
    fill = shape.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = BG_TOP
    stops[0].position = 0.0
    stops[1].color.rgb = BG_BOTTOM
    stops[1].position = 1.0
    fill.gradient_angle = 115
    shape.element.getparent().remove(shape.element)
    slide.shapes._spTree.insert(2, shape.element)
    return shape


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_w=Pt(1), shadow=False, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
    shape.shadow.inherit = False
    if radius:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, l, t, w, h, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.05, wrap=True):
    tb = slide.shapes.add_textbox(Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color
    return tb


def add_rich_text(slide, l, t, w, h, runs, size=14, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                   font=FONT, bold=False, line_spacing=1.05):
    """runs: list of (text, color, bold_override) 拼接一段带高亮的文字"""
    tb = slide.shapes.add_textbox(Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for seg in runs:
        seg_text, seg_color = seg[0], seg[1]
        seg_bold = seg[2] if len(seg) > 2 else bold
        r = p.add_run()
        r.text = seg_text
        r.font.size = Pt(size)
        r.font.bold = seg_bold
        r.font.name = font
        r.font.color.rgb = seg_color
    return tb


def add_header(slide, prs, title_text, title_size=23):
    """左上角项目背景徽标 + 居中大标题"""
    badge = add_rect(slide, Emu(int(0.35 * 914400)), Emu(int(0.22 * 914400)),
                      Emu(int(2.0 * 914400)), Emu(int(0.5 * 914400)), fill_color=BADGE_FILL, radius=0.5)
    tf = badge.text_frame
    tf.margin_left = Emu(int(0.05 * 914400))
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "‹‹ 项 目 背 景"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = WHITE
    # 标题
    add_text(slide, Emu(int(2.6 * 914400)), Emu(int(0.2 * 914400)),
              prs.slide_width - Emu(int(5.2 * 914400)), Emu(int(0.6 * 914400)),
              title_text, size=title_size, color=GOLD, bold=True, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
    # 发丝分隔线
    line = slide.shapes.add_connector(1, Emu(int(0.35 * 914400)), Emu(int(0.95 * 914400)),
                                       prs.slide_width - Emu(int(0.35 * 914400)), Emu(int(0.95 * 914400)))
    line.line.color.rgb = PANEL_LINE
    line.line.width = Pt(0.75)


def add_bottom_bar(slide, prs, runs, y=None):
    h = Emu(int(0.62 * 914400))
    y = y if y else prs.slide_height - h - Emu(int(0.18 * 914400))
    add_rect(slide, Emu(int(0.35 * 914400)), y, prs.slide_width - Emu(int(0.7 * 914400)), h,
              fill_color=BOTTOM_BAR_FILL, line_color=GOLD, line_w=Pt(1), radius=0.25)
    add_rich_text(slide, Emu(int(0.6 * 914400)), y, prs.slide_width - Emu(int(1.2 * 914400)), h,
                  runs, size=15, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)


def add_panel(slide, l, t, w, h, header=None, header_color=CYAN):
    panel = add_rect(slide, l, t, w, h, fill_color=PANEL_FILL, line_color=PANEL_LINE, line_w=Pt(1), radius=0.04)
    if header:
        add_text(slide, l + Emu(int(0.12 * 914400)), t + Emu(int(0.06 * 914400)),
                  w - Emu(int(0.24 * 914400)), Emu(int(0.35 * 914400)), header,
                  size=15, color=header_color, bold=True, align=PP_ALIGN.LEFT)
    return panel


def add_picture_framed(slide, path, l, t, w, h, line_color=PANEL_LINE):
    pic = slide.shapes.add_picture(path, Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
    pic.line.color.rgb = line_color
    pic.line.width = Pt(0.75)
    return pic


def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_gradient_bg(slide, prs)
    return slide


IN = 914400  # EMU per inch
