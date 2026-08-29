#!/usr/bin/env python3
"""从 PPTX 按页提取所有图片，按版面坐标命名导出。"""

import json
from pathlib import Path
from pptx import Presentation
from PIL import Image
import io
import sys

def extract_pptx_images(pptx_path: str, output_dir: str):
    prs = Presentation(pptx_path)
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    
    manifest = []
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        img_idx = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                img_idx += 1
                image = shape.image
                ext = image.content_type.split("/")[-1].lower()
                if ext == "jpeg":
                    ext = "jpg"
                elif ext not in ["jpg", "png", "webp", "gif", "bmp"]:
                    ext = "png"
                
                # 版面坐标（EMU → cm）
                left_cm = shape.left / 914400 * 2.54
                top_cm = shape.top / 914400 * 2.54
                width_cm = shape.width / 914400 * 2.54
                height_cm = shape.height / 914400 * 2.54
                
                # 命名规则：P{页码}_{序号}_L{left}T{top}W{width}H{height}.{ext}
                fname = f"P{slide_idx:02d}_{img_idx:02d}_L{left_cm:.1f}T{top_cm:.1f}W{width_cm:.1f}H{height_cm:.1f}.{ext}"
                fpath = out / fname
                
                with open(fpath, "wb") as f:
                    f.write(image.blob)
                
                # 读取实际像素尺寸
                try:
                    pil_img = Image.open(io.BytesIO(image.blob))
                    px_w, px_h = pil_img.size
                except Exception:
                    px_w, px_h = 0, 0
                
                manifest.append({
                    "page": slide_idx,
                    "index": img_idx,
                    "filename": fname,
                    "rel_path": str(fpath),
                    "left_cm": round(left_cm, 1),
                    "top_cm": round(top_cm, 1),
                    "width_cm": round(width_cm, 1),
                    "height_cm": round(height_cm, 1),
                    "pixel_width": px_w,
                    "pixel_height": px_h,
                    "aspect_ratio": f"{px_w}:{px_h}",
                })
    
    # 输出清单
    manifest_path = out / "image_manifest.json"
    with open(manifest_path, "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)
    
    print(f"导出 {len(manifest)} 张图片到 {output_dir}/")
    print(f"清单: {manifest_path}")
    return manifest

if __name__ == "__main__":
    pptx_file = sys.argv[1] if len(sys.argv) > 1 else "毕业答辩.pptx"
    out_dir = sys.argv[2] if len(sys.argv) > 2 else "03-素材库/毕业答辩素材"
    extract_pptx_images(pptx_file, out_dir)
